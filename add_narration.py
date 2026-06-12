"""
Embed narration MP3s into the existing Rexora-Systems-Pitch.pptx so each
audio file plays automatically when its slide appears, and each slide
advances when narration finishes (+0.6s buffer).

Inputs:  Rexora-Systems-Pitch.pptx  +  audio/slide_*.mp3
Output:  Rexora-Systems-Pitch-Narrated.pptx
"""

import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree
from mutagen.mp3 import MP3

SRC = "Rexora-Systems-Pitch.pptx"
DST = "Rexora-Systems-Pitch-Narrated.pptx"
AUDIO_DIR = "audio"
END_BUFFER_SEC = 0.6

NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def mp3_duration_sec(path: str) -> float:
    return MP3(path).info.length


def embed_autoplay_audio(slide, audio_path: str) -> int:
    """Add audio shape (hidden offscreen) and inject autoplay timing.
    Returns the shape id used for the media."""
    # add_movie supports audio mime
    pic = slide.shapes.add_movie(
        audio_path,
        left=Emu(-914400),                # off-slide so icon isn't visible
        top=Emu(-914400),
        width=Inches(0.3),
        height=Inches(0.3),
        mime_type='audio/mpeg',
    )

    # Discover shape id
    sp = pic._element                          # p:pic
    cNvPr = sp.find('.//' + qn('p:cNvPr'))
    if cNvPr is None:
        cNvPr = sp.find('.//' + qn('p:nvPicPr') + '/' + qn('p:cNvPr'))
    sp_id = cNvPr.get('id')

    # Build timing XML for autoplay
    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="1" presetClass="mediacall"
                                     presetSubtype="0" fill="hold" nodeType="afterEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="indefinite" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
          <p:audio>
            <p:cMediaNode vol="80000">
              <p:cTn id="7" fill="hold" display="0" masterRel="sameClick">
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
              </p:cTn>
              <p:tgtEl><p:spTgt spid="{sp_id}"/></p:tgtEl>
            </p:cMediaNode>
          </p:audio>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''
    new_timing = etree.fromstring(timing_xml)

    sld = slide.element  # p:sld root
    existing = sld.find(qn('p:timing'))
    if existing is not None:
        sld.remove(existing)
    sld.append(new_timing)

    return int(sp_id)


def set_slide_advance(slide, seconds: float) -> None:
    """Set auto-advance time on the slide's transition (creates one if missing)."""
    sld = slide.element
    transition = sld.find(qn('p:transition'))
    if transition is None:
        transition = etree.SubElement(sld, qn('p:transition'))
        # transition must come before p:timing
        timing = sld.find(qn('p:timing'))
        if timing is not None:
            sld.remove(transition)
            sld.insert(list(sld).index(timing), transition)
        fade = etree.SubElement(transition, qn('p:fade'))
    transition.set('spd', 'med')
    transition.set('advTm', str(int(seconds * 1000)))
    transition.set('advClick', '0')


def main() -> None:
    if not os.path.exists(SRC):
        raise SystemExit(f"Missing {SRC}. Run build_pitch_deck.py first.")

    prs = Presentation(SRC)
    print(f"Loaded {SRC}  ({len(prs.slides)} slides)")

    for i, slide in enumerate(prs.slides, start=1):
        audio_path = os.path.join(AUDIO_DIR, f"slide_{i}.mp3")
        if not os.path.exists(audio_path):
            print(f"  slide {i}: no audio, skipped")
            continue
        dur = mp3_duration_sec(audio_path)
        embed_autoplay_audio(slide, audio_path)
        set_slide_advance(slide, dur + END_BUFFER_SEC)
        print(f"  slide {i}: {dur:5.2f}s narration  (advance {dur+END_BUFFER_SEC:.2f}s)")

    prs.save(DST)
    print(f"\nSaved {DST}")


if __name__ == "__main__":
    main()
