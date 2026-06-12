"""
Build the Rexora Systems pitch deck (PowerPoint).
Designed for the African SMEs TikTok script. 16:9 widescreen,
brand-aligned dark theme, slides auto-advance so the file can be
exported to MP4 directly via PowerPoint > File > Export > Create a Video.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Brand tokens ----------
INK         = RGBColor(0x0A, 0x0A, 0x0F)
SUBINK      = RGBColor(0x15, 0x15, 0x1C)
BRAND_GREEN = RGBColor(0x4A, 0xDE, 0x80)
BRAND_BLUE  = RGBColor(0x60, 0xA5, 0xFA)
BRAND_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300    = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400    = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500    = RGBColor(0x6B, 0x72, 0x80)

DISPLAY_FONT = "Poppins"
BODY_FONT    = "Inter"

# ---------- Presentation ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line=False, transparency=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if not line:
        shp.line.fill.background()
    if transparency is not None:
        # set fill transparency via XML
        sp = shp.fill._xPr
        solidFill = sp.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', str(int((1 - transparency) * 100000)))
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, w, h, fill_color, transparency=0.7):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    sp = shp.fill._xPr
    solidFill = sp.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(int((1 - transparency) * 100000)))
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             font=BODY_FONT, size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs = [{'text','font','size','bold','italic','color'}, ...]
       New paragraph signaled by run with key 'br': True."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in runs:
        if r.get('br'):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            continue
        run = p.add_run()
        run.text = r.get('text', '')
        run.font.name = r.get('font', BODY_FONT)
        run.font.size = Pt(r.get('size', 18))
        run.font.bold = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
    return tb


def add_logo_mark(slide, x, y, size=Inches(0.55)):
    """Tri-color triangular prism mark — composed of 3 triangles."""
    # triangle 1 — top diamond (green) — points: top, mid-right, mid, mid-left
    # We'll use freeform shapes
    s = size
    # green diamond top
    g = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x + s * 0.31, y, s * 0.38, s * 0.6)
    g.fill.solid(); g.fill.fore_color.rgb = BRAND_GREEN
    g.line.fill.background(); g.shadow.inherit = False

    # blue triangle bottom-left
    b = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, x, y + s * 0.35, s * 0.55, s * 0.6)
    b.fill.solid(); b.fill.fore_color.rgb = BRAND_BLUE
    b.line.fill.background(); b.shadow.inherit = False

    # purple triangle bottom-right (flipped)
    pu = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, x + s * 0.45, y + s * 0.35, s * 0.55, s * 0.6)
    pu.rotation = 270
    pu.fill.solid(); pu.fill.fore_color.rgb = BRAND_PURPLE
    pu.line.fill.background(); pu.shadow.inherit = False


def add_brand_header(slide):
    add_logo_mark(slide, Inches(0.5), Inches(0.4))
    add_runs(slide, Inches(1.15), Inches(0.42), Inches(5), Inches(0.5),
             [
                 {'text': 'REXORA ', 'font': DISPLAY_FONT, 'size': 14, 'bold': True, 'color': WHITE},
                 {'text': 'SYSTEMS', 'font': DISPLAY_FONT, 'size': 9, 'color': BRAND_PURPLE},
             ])


def set_auto_advance(slide, seconds):
    """Auto-advance slide after N seconds (no click required)."""
    spTree = slide.element
    transition = etree.SubElement(spTree, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advTm', str(int(seconds * 1000)))
    # also disable click-advance
    transition.set('advClick', '0')
    fade = etree.SubElement(transition, qn('p:fade'))


def add_glow_orb(slide, cx, cy, radius, color, transparency=0.85):
    add_oval(slide, cx - radius, cy - radius, radius * 2, radius * 2, color, transparency)


# ============================================================
#                       SLIDE BUILDERS
# ============================================================

def build_slide_1_hook():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_glow_orb(s, Inches(2), Inches(2), Inches(3.2), BRAND_GREEN, 0.82)
    add_glow_orb(s, Inches(11), Inches(6), Inches(3.5), BRAND_PURPLE, 0.82)
    add_glow_orb(s, Inches(7), Inches(4), Inches(3), BRAND_BLUE, 0.85)

    add_brand_header(s)

    # Eyebrow
    add_text(s, Inches(1), Inches(2.3), Inches(11.3), Inches(0.4),
             "FOR AFRICAN SMES",
             font=BODY_FONT, size=14, bold=True, color=BRAND_GREEN,
             align=PP_ALIGN.CENTER)

    # Hook
    add_text(s, Inches(0.7), Inches(2.85), Inches(12), Inches(2.5),
             "If your business runs on\nWhatsApp chats and a notebook…",
             font=DISPLAY_FONT, size=46, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    add_text(s, Inches(1), Inches(5.4), Inches(11.3), Inches(0.6),
             "this is for you.",
             font=DISPLAY_FONT, size=32, italic=True, color=BRAND_GREEN,
             align=PP_ALIGN.CENTER)

    # Footer brand
    add_text(s, Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
             "Rexora Systems  ·  Digital Systems. Smarter Operations. Bigger Impact.",
             font=BODY_FONT, size=11, color=GRAY_500,
             align=PP_ALIGN.CENTER)

    set_auto_advance(s, 4)


def build_slide_2_reality():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_glow_orb(s, Inches(11), Inches(2), Inches(3), BRAND_BLUE, 0.85)
    add_brand_header(s)

    # eyebrow
    add_text(s, Inches(1), Inches(1.8), Inches(11.3), Inches(0.4),
             "THE REALITY",
             font=BODY_FONT, size=14, bold=True, color=GRAY_400,
             align=PP_ALIGN.CENTER)

    # Big stat
    add_runs(s, Inches(1), Inches(2.4), Inches(11.3), Inches(2.5),
             [
                 {'text': '9 out of 10', 'font': DISPLAY_FONT, 'size': 100, 'bold': True, 'color': BRAND_GREEN},
             ],
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.7),
             "African SMEs we audit have the same problem.",
             font=DISPLAY_FONT, size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(2), Inches(5.6), Inches(9.3), Inches(1.4),
             "Every booking, every payment, every customer\nlives in someone's head — or a chat thread.",
             font=BODY_FONT, size=18, italic=True, color=GRAY_300,
             align=PP_ALIGN.CENTER, line_spacing=1.4)

    set_auto_advance(s, 5)


def build_slide_3_today():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_brand_header(s)

    add_text(s, Inches(1), Inches(1.3), Inches(11.3), Inches(0.5),
             "TODAY, YOUR BUSINESS RUNS ON",
             font=BODY_FONT, size=14, bold=True, color=GRAY_400,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(1), Inches(1.85), Inches(11.3), Inches(0.9),
             "A patchwork, not a system.",
             font=DISPLAY_FONT, size=38, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # 4-card grid
    items = [
        ("WhatsApp",        "for bookings",                BRAND_GREEN),
        ("A notebook",      "for payments",                BRAND_BLUE),
        ("One person's mind","for everything else",        BRAND_PURPLE),
        ("An IG bio",       "instead of a website",        BRAND_GREEN),
    ]
    card_w = Inches(2.7)
    card_h = Inches(2.6)
    gap    = Inches(0.25)
    total  = card_w * 4 + gap * 3
    start_x = (SW - total) // 2
    y      = Inches(3.4)

    for i, (title, sub, color) in enumerate(items):
        x = start_x + (card_w + gap) * i
        # card bg
        card = add_rect(s, x, y, card_w, card_h, SUBINK)
        # accent bar
        add_rect(s, x, y, card_w, Inches(0.08), color)
        # number
        add_text(s, x + Inches(0.25), y + Inches(0.3), card_w, Inches(0.45),
                 f"0{i+1}",
                 font=DISPLAY_FONT, size=14, bold=True, color=color)
        # title
        add_text(s, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.9),
                 title,
                 font=DISPLAY_FONT, size=20, bold=True, color=WHITE, line_spacing=1.1)
        # sub
        add_text(s, x + Inches(0.25), y + Inches(1.85), card_w - Inches(0.5), Inches(0.6),
                 sub,
                 font=BODY_FONT, size=13, color=GRAY_400)

    set_auto_advance(s, 5)


def build_slide_4_cost():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_glow_orb(s, Inches(2), Inches(7), Inches(3), BRAND_PURPLE, 0.88)
    add_brand_header(s)

    add_text(s, Inches(1), Inches(1.3), Inches(11.3), Inches(0.5),
             "THE HIDDEN COST",
             font=BODY_FONT, size=14, bold=True, color=BRAND_PURPLE,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(1), Inches(1.85), Inches(11.3), Inches(0.9),
             "When there's no system…",
             font=DISPLAY_FONT, size=38, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    rows = [
        ("That person travels",      "→  business stops."),
        ("The chat gets archived",   "→  money is lost."),
        ("You want to grow",         "→  you can't. There's no system to grow into."),
    ]
    y = Inches(3.4)
    for i, (cause, effect) in enumerate(rows):
        ry = y + Inches(0.95 * i)
        # row card
        add_rect(s, Inches(1.5), ry, Inches(10.3), Inches(0.75), SUBINK)
        # left - cause
        add_text(s, Inches(1.85), ry + Inches(0.18), Inches(4.5), Inches(0.45),
                 cause,
                 font=DISPLAY_FONT, size=18, bold=True, color=WHITE)
        # right - effect
        add_text(s, Inches(6.2), ry + Inches(0.18), Inches(5.5), Inches(0.45),
                 effect,
                 font=BODY_FONT, size=17, color=BRAND_GREEN, italic=True)

    set_auto_advance(s, 6)


def build_slide_5_shift():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_glow_orb(s, Inches(6.66), Inches(3.75), Inches(5), BRAND_GREEN, 0.85)
    add_glow_orb(s, Inches(2),    Inches(2),    Inches(3), BRAND_BLUE, 0.88)
    add_glow_orb(s, Inches(11),   Inches(6),    Inches(3), BRAND_PURPLE, 0.88)

    add_brand_header(s)

    add_runs(s, Inches(0.7), Inches(2.6), Inches(12), Inches(3),
             [
                 {'text': 'Your business deserves a ', 'font': DISPLAY_FONT, 'size': 50, 'bold': True, 'color': WHITE},
                 {'text': 'system', 'font': DISPLAY_FONT, 'size': 50, 'bold': True, 'color': BRAND_GREEN},
                 {'text': ',', 'font': DISPLAY_FONT, 'size': 50, 'bold': True, 'color': WHITE},
                 {'br': True},
                 {'text': 'not a notebook.', 'font': DISPLAY_FONT, 'size': 50, 'bold': True, 'color': WHITE},
             ],
             align=PP_ALIGN.CENTER, line_spacing=1.15)

    set_auto_advance(s, 4)


def build_slide_6_pillars():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_brand_header(s)

    add_text(s, Inches(1), Inches(1.3), Inches(11.3), Inches(0.5),
             "WHAT WE BUILD",
             font=BODY_FONT, size=14, bold=True, color=GRAY_400,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(1), Inches(1.85), Inches(11.3), Inches(0.9),
             "Three outcomes. Zero fluff.",
             font=DISPLAY_FONT, size=38, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    pillars = [
        ("DRIVE GROWTH",  "Websites & funnels that sell while you sleep.",       BRAND_GREEN),
        ("AUTOMATE",      "Bookings, payments & reminders that run themselves.", BRAND_BLUE),
        ("ELEVATE BRAND", "A premium presence that builds trust.",               BRAND_PURPLE),
    ]
    card_w = Inches(3.7)
    card_h = Inches(3)
    gap    = Inches(0.35)
    total  = card_w * 3 + gap * 2
    start_x = (SW - total) // 2
    y      = Inches(3.3)

    for i, (title, sub, color) in enumerate(pillars):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, card_h, SUBINK)
        add_rect(s, x, y, card_w, Inches(0.1), color)

        # icon dot
        dot = slide_icon = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.5),
                                               Inches(0.45), Inches(0.45))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background(); dot.shadow.inherit = False

        add_text(s, x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.6), Inches(0.5),
                 title,
                 font=DISPLAY_FONT, size=18, bold=True, color=color)
        add_text(s, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(1.2),
                 sub,
                 font=BODY_FONT, size=15, color=GRAY_300, line_spacing=1.4)

    add_text(s, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
             "Built for African SMEs — schools, churches, clinics, and growing businesses.",
             font=BODY_FONT, size=13, italic=True, color=GRAY_500,
             align=PP_ALIGN.CENTER)

    set_auto_advance(s, 6)


def build_slide_7_proof():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_glow_orb(s, Inches(11), Inches(6), Inches(3), BRAND_GREEN, 0.88)
    add_brand_header(s)

    add_text(s, Inches(1), Inches(1.3), Inches(11.3), Inches(0.5),
             "PROOF",
             font=BODY_FONT, size=14, bold=True, color=BRAND_GREEN,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(1), Inches(1.85), Inches(11.3), Inches(0.9),
             "Trusted by ambitious teams across Africa.",
             font=DISPLAY_FONT, size=32, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # logo strip (text only since we won't load image files)
    brands = ["ARTEESANS", "THE SQUARCLE", "CJB MINISTRY", "VENDORA", "LAMAH HUB", "DREAMPAWS"]
    strip_y = Inches(3.0)
    add_rect(s, Inches(0.7), strip_y, Inches(11.93), Inches(0.85), SUBINK)
    col_w = Inches(11.93) / len(brands)
    for i, b in enumerate(brands):
        add_text(s, Inches(0.7) + col_w * i, strip_y + Inches(0.25),
                 col_w, Inches(0.5),
                 b,
                 font=DISPLAY_FONT, size=12, bold=True, color=GRAY_400,
                 align=PP_ALIGN.CENTER)

    # quote card
    quote_y = Inches(4.4)
    add_rect(s, Inches(1.5), quote_y, Inches(10.3), Inches(2.3), SUBINK)
    add_rect(s, Inches(1.5), quote_y, Inches(0.1), Inches(2.3), BRAND_GREEN)

    add_text(s, Inches(1.85), quote_y + Inches(0.3), Inches(9.6), Inches(0.4),
             "★★★★★",
             font=BODY_FONT, size=14, color=BRAND_GREEN)
    add_text(s, Inches(1.85), quote_y + Inches(0.7), Inches(9.6), Inches(1.1),
             "“Bookings doubled overnight. Clients pay before they walk in —\nit's like having a salesperson who never sleeps.”",
             font=DISPLAY_FONT, size=18, italic=True, color=WHITE, line_spacing=1.35)
    add_text(s, Inches(1.85), quote_y + Inches(1.85), Inches(9.6), Inches(0.4),
             "— Vendora  ·  Business",
             font=BODY_FONT, size=12, color=GRAY_400)

    set_auto_advance(s, 6)


def build_slide_8_cta():
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_glow_orb(s, Inches(6.66), Inches(3.75), Inches(5), BRAND_GREEN, 0.82)
    add_glow_orb(s, Inches(2),    Inches(6),    Inches(3), BRAND_PURPLE, 0.88)
    add_brand_header(s)

    add_text(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.5),
             "FREE BUSINESS AUDIT",
             font=BODY_FONT, size=14, bold=True, color=BRAND_GREEN,
             align=PP_ALIGN.CENTER)

    add_runs(s, Inches(0.7), Inches(2.3), Inches(12), Inches(2),
             [
                 {'text': "Let's build something that ", 'font': DISPLAY_FONT, 'size': 46, 'bold': True, 'color': WHITE},
                 {'br': True},
                 {'text': "makes you money", 'font': DISPLAY_FONT, 'size': 46, 'bold': True, 'color': BRAND_GREEN},
                 {'text': ".", 'font': DISPLAY_FONT, 'size': 46, 'bold': True, 'color': WHITE},
             ],
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    add_text(s, Inches(2), Inches(4.3), Inches(9.3), Inches(0.7),
             "We'll map exactly what to automate, what to launch, and what it'll be worth to your bottom line.",
             font=BODY_FONT, size=15, color=GRAY_300,
             align=PP_ALIGN.CENTER, line_spacing=1.4)

    # CTA pill
    pill_w = Inches(5.2)
    pill_h = Inches(0.85)
    pill_x = (SW - pill_w) // 2
    pill_y = Inches(5.35)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_x, pill_y, pill_w, pill_h)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = BRAND_GREEN
    pill.line.fill.background(); pill.shadow.inherit = False
    add_text(s, pill_x, pill_y + Inches(0.18), pill_w, Inches(0.55),
             "Get a free system audit  →",
             font=DISPLAY_FONT, size=20, bold=True, color=INK,
             align=PP_ALIGN.CENTER)

    # Contact row
    add_runs(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             [
                 {'text': 'rexorasystems.com', 'font': BODY_FONT, 'size': 14, 'bold': True, 'color': WHITE},
                 {'text': '   ·   ', 'font': BODY_FONT, 'size': 14, 'color': GRAY_500},
                 {'text': 'hello@rexorasystems.com', 'font': BODY_FONT, 'size': 14, 'bold': True, 'color': WHITE},
             ],
             align=PP_ALIGN.CENTER)

    set_auto_advance(s, 6)


# ---------- Build ----------
build_slide_1_hook()
build_slide_2_reality()
build_slide_3_today()
build_slide_4_cost()
build_slide_5_shift()
build_slide_6_pillars()
build_slide_7_proof()
build_slide_8_cta()

OUT = "Rexora-Systems-Pitch.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
